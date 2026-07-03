"""Ingest step: a LOADER REGISTRY that turns any supported document into the
common record shape {loc, order, text[, ocr]}.

Digital formats load their text directly. Anything with no extractable text
(scanned PDF pages, image files, handwriting) falls through to the adaptive
OCR engine (src/ocr.py) and is tagged ocr=True so the answer step can add a
"check the equation" nudge when needed.
"""
import csv
import io
import re
from pathlib import Path

import pymupdf
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook
from PIL import Image

from src.config import DOCS_DIR, OCR_DPI, OCR_MIN_TEXT


def clean(text: str) -> str:
    text = text.replace("\x00", " ")
    return re.sub(r"\s+", " ", text).strip()


# ---------- digital loaders ----------
def _load_pdf(path: Path) -> list[dict]:
    """Per page: use embedded text if present, else OCR the rendered page."""
    doc = pymupdf.open(str(path))
    out = []
    for i, page in enumerate(doc, start=1):
        t = clean(page.get_text() or "")
        if len(t) >= OCR_MIN_TEXT:                       # real digital text
            out.append({"loc": f"p{i}", "order": i, "text": t})
        else:                                            # scanned / handwritten
            from src.ocr import ocr_image                # lazy: only import API here
            png = page.get_pixmap(dpi=OCR_DPI).tobytes("png")
            ocr_t = clean(ocr_image(png))
            if ocr_t:
                out.append({"loc": f"p{i}", "order": i, "text": ocr_t, "ocr": True})
    return out


def _load_docx(path: Path) -> list[dict]:
    doc = Document(str(path))
    out, buf, order, sec = [], [], 0, 0
    label = "body"

    def flush():
        nonlocal order
        if buf:
            t = clean(" ".join(buf))
            if t:
                order += 1
                out.append({"loc": label, "order": order, "text": t})

    for p in doc.paragraphs:
        style = (p.style.name or "").lower()
        if style.startswith("heading") or style.startswith("title"):
            flush(); buf.clear(); sec += 1
            label = clean(p.text)[:40] or f"section {sec}"
        elif p.text.strip():
            buf.append(p.text)
    flush()
    for tbl in doc.tables:
        for row in tbl.rows:
            line = clean(" | ".join(c.text for c in row.cells))
            if line:
                order += 1
                out.append({"loc": "table", "order": order, "text": line})
    return out


def _load_pptx(path: Path) -> list[dict]:
    out = []
    for i, slide in enumerate(Presentation(str(path)).slides, start=1):
        parts = []
        for sh in slide.shapes:
            if sh.has_text_frame:
                for para in sh.text_frame.paragraphs:
                    line = "".join(r.text for r in para.runs)
                    if line.strip():
                        parts.append(line)
            if sh.has_table:
                for row in sh.table.rows:
                    parts.append(" | ".join(c.text for c in row.cells))
        t = clean("\n".join(parts))
        if t:
            out.append({"loc": f"slide {i}", "order": i, "text": t})
    return out


def _load_text(path: Path) -> list[dict]:
    t = clean(path.read_text(encoding="utf-8", errors="ignore"))
    return [{"loc": "text", "order": 1, "text": t}] if t else []


def _load_csv(path: Path) -> list[dict]:
    out = []
    with open(path, newline="", encoding="utf-8", errors="ignore") as f:
        rows = list(csv.reader(f))
    if not rows:
        return out
    header = rows[0]
    for i, row in enumerate(rows[1:], start=1):
        line = "; ".join(f"{h}: {v}" for h, v in zip(header, row) if str(v).strip())
        if line:
            out.append({"loc": f"row {i}", "order": i, "text": line})
    return out


def _load_xlsx(path: Path) -> list[dict]:
    wb = load_workbook(str(path), read_only=True, data_only=True)
    out, order = [], 0
    for ws in wb.worksheets:
        rows = list(ws.iter_rows(values_only=True))
        if not rows:
            continue
        header = [str(c) if c is not None else "" for c in rows[0]]
        for r, row in enumerate(rows[1:], start=1):
            line = "; ".join(f"{h}: {v}" for h, v in zip(header, row) if v not in (None, ""))
            if line:
                order += 1
                out.append({"loc": f"{ws.title} row {r}", "order": order, "text": line})
    return out


# ---------- image loader (pure OCR) ----------
def _load_image(path: Path) -> list[dict]:
    from src.ocr import ocr_image
    img = Image.open(path).convert("RGB")
    buf = io.BytesIO(); img.save(buf, format="PNG")
    text = clean(ocr_image(buf.getvalue()))
    return [{"loc": "image", "order": 1, "text": text, "ocr": True}] if text else []


IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".webp", ".tiff", ".bmp"}
LOADERS = {
    ".pdf": _load_pdf,
    ".docx": _load_docx,
    ".pptx": _load_pptx,
    ".txt": _load_text,
    ".md": _load_text,
    ".csv": _load_csv,
    ".xlsx": _load_xlsx,
    **{e: _load_image for e in IMAGE_EXTS},
}
SUPPORTED = set(LOADERS)


def list_documents() -> list[Path]:
    return sorted(p for p in DOCS_DIR.iterdir()
                  if p.is_file() and p.suffix.lower() in SUPPORTED)


def unsupported_files() -> list[Path]:
    return sorted(p for p in DOCS_DIR.iterdir()
                  if p.is_file() and p.suffix.lower() not in SUPPORTED
                  and not p.name.startswith("."))


def load_document(path: Path) -> list[dict]:
    return LOADERS[path.suffix.lower()](path)


def source_label(path: Path) -> str:
    return re.sub(r"[\s_]+", " ", path.stem).strip()


if __name__ == "__main__":
    for d in list_documents():
        recs = load_document(d)
        n_ocr = sum(1 for r in recs if r.get("ocr"))
        print(f"{source_label(d)} ({d.suffix}) -> {len(recs)} records"
              f"{f'  [{n_ocr} via OCR]' if n_ocr else ''}")
